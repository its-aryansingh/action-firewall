"""Generate docs/Razorpay_Buildathon_Action_Firewall_Deck.pptx matching pitch-deck.html."""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 16:9 widescreen dimensions
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Color palette matching pitch-deck.html
C_INK = RGBColor(11, 16, 32)        # #0b1020
C_PANEL = RGBColor(21, 30, 58)      # #151e3a
C_PANEL2 = RGBColor(28, 41, 75)     # #1c294b
C_EDGE = RGBColor(43, 58, 103)      # #2b3a67
C_BLUE = RGBColor(51, 149, 255)     # #3395ff
C_GREEN = RGBColor(35, 214, 111)    # #23d66f
C_RED = RGBColor(255, 63, 104)      # #ff3f68
C_AMBER = RGBColor(255, 189, 46)    # #ffbd2e
C_TEXT = RGBColor(245, 247, 255)    # #f5f7ff
C_MUTED = RGBColor(156, 174, 208)   # #9caed0


def set_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = C_INK


def add_header(slide, kicker: str, title: str, subtitle: str = ""):
    # Kicker
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = kicker.upper()
    p.font.name = "Calibri"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = C_BLUE

    # Title
    tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.9), Inches(11.7), Inches(0.8))
    tf_title = tb_title.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
    p_t = tf_title.paragraphs[0]
    p_t.text = title
    p_t.font.name = "Georgia"
    p_t.font.size = Pt(28)
    p_t.font.bold = True
    p_t.font.color.rgb = C_TEXT

    if subtitle:
        tb_sub = slide.shapes.add_textbox(Inches(0.8), Inches(1.75), Inches(11.7), Inches(0.5))
        tf_sub = tb_sub.text_frame
        tf_sub.word_wrap = True
        tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
        p_s = tf_sub.paragraphs[0]
        p_s.text = subtitle
        p_s.font.name = "Calibri"
        p_s.font.size = Pt(14)
        p_s.font.color.rgb = C_MUTED


def add_footer(slide, left_text: str, right_text: str = "Aryan Singh"):
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = left_text
    p.font.name = "Consolas"
    p.font.size = Pt(10)
    p.font.color.rgb = C_MUTED

    tb_r = slide.shapes.add_textbox(Inches(8.0), Inches(6.8), Inches(4.5), Inches(0.3))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_top = tf_r.margin_right = tf_r.margin_bottom = 0
    p_r = tf_r.paragraphs[0]
    p_r.text = right_text
    p_r.alignment = PP_ALIGN.RIGHT
    p_r.font.name = "Calibri"
    p_r.font.size = Pt(10)
    p_r.font.color.rgb = C_MUTED


def add_card(slide, left, top, width, height, title="", body="", title_color=C_TEXT, border_color=C_EDGE, bg_color=C_PANEL):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_top = Inches(0.2)
    tf.margin_right = Inches(0.25)
    tf.margin_bottom = Inches(0.2)

    if title:
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.name = "Georgia" if len(title) > 20 else "Calibri"
        p_title.font.size = Pt(16)
        p_title.font.bold = True
        p_title.font.color.rgb = title_color
        p_title.space_after = Pt(8)

    if body:
        p_body = tf.add_paragraph() if title else tf.paragraphs[0]
        p_body.text = body
        p_body.font.name = "Calibri"
        p_body.font.size = Pt(12)
        p_body.font.color.rgb = C_MUTED
    return shape


def build_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    # --- SLIDE 1: Title ---
    s1 = prs.slides.add_slide(blank_layout)
    set_bg(s1)
    add_header(s1, "Razorpay AI Buildathon 2026 · Track 01", "Action Firewall — Safe Autopilot Checkout", "One approval for the job. Zero authority beyond it.")

    # Pills / Highlights
    add_card(s1, Inches(0.8), Inches(2.6), Inches(5.6), Inches(1.2),
             "AI DRAFTS & RECOVERS", "Zero autonomous payment authority. Translates intent and ranks eligible catalog items.",
             title_color=C_GREEN, border_color=C_GREEN)
    add_card(s1, Inches(6.8), Inches(2.6), Inches(5.7), Inches(1.2),
             "PURCHASE ENVELOPE", "One exact Action Grant · ₹417 / ₹600. Versioned, hash-bound spend fence.",
             title_color=C_BLUE, border_color=C_BLUE)

    # Lead box
    add_card(s1, Inches(0.8), Inches(4.2), Inches(11.7), Inches(1.8),
             "Agentic Commerce Control Plane",
             "An AI may plan and recover a purchase. Only a deterministic, versioned Purchase Envelope can authorize the final Razorpay action.\n\n"
             "• Chat proposes; actuator calls registered create_payment_link once\n"
             "• Re-verified inside SQLite BEGIN IMMEDIATE before any network I/O\n"
             "• Dual HMAC-SHA256 Action Receipt survives grant settlement",
             title_color=C_TEXT)

    add_footer(s1, "Razorpay MCP-compatible actuator · deterministic policy · 102 tests · offline demo", "Aryan Singh · Track 01")

    # --- SLIDE 2: The Control Gap ---
    s2 = prs.slides.add_slide(blank_layout)
    set_bg(s2)
    add_header(s2, "The control gap", "A model is not a payment principal.",
               "Natural language expresses intent. It cannot carry durable authority across catalog substitutions, merchant changes, retries, and provider timeouts.")

    add_card(s2, Inches(0.8), Inches(2.4), Inches(3.7), Inches(2.6),
             "1 · Prompt & Plan Drift",
             "Stock loss or conversational ambiguity can nudge an agent toward unapproved items, categories, or price spikes without human realization.",
             title_color=C_RED, border_color=C_RED)
    add_card(s2, Inches(4.8), Inches(2.4), Inches(3.7), Inches(2.6),
             "2 · Merchant & State Drift",
             "An item sells out and the agent re-routes to an unapproved seller. A general permission becomes an uninspected blank check.",
             title_color=C_AMBER, border_color=C_AMBER)
    add_card(s2, Inches(8.8), Inches(2.4), Inches(3.7), Inches(2.6),
             "3 · Distributed Ambiguity",
             "Network timeouts and retries duplicate external payment actions unless exposure is tracked atomically before dispatch.",
             title_color=C_BLUE, border_color=C_BLUE)

    add_card(s2, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.1),
             "The missing primitive: Purchase Envelope",
             "Bounded human approval, exact machine verification. Track 01 requires explainable, bounded and gated money actions with full auditability.",
             title_color=C_TEXT, bg_color=C_PANEL2)
    add_footer(s2, "First-party requirement · Razorpay AI Buildathon 2026")

    # --- SLIDE 3: Architecture ---
    s3 = prs.slides.add_slide(blank_layout)
    set_bg(s3)
    add_header(s3, "Architecture", "AI plans and recovers. Deterministic code gates money.",
               "The model drafts intent; one shopper click activates it; server-side code verifies exact catalog facts before any provider call.")

    add_card(s3, Inches(0.8), Inches(2.4), Inches(2.7), Inches(2.0),
             "1 · Draft Intent", "AI or deterministic fallback drafts required slots and budget from shopping goal.", title_color=C_BLUE)
    add_card(s3, Inches(3.8), Inches(2.4), Inches(2.7), Inches(2.0),
             "2 · Activate Once", "Shopper reviews budget, merchant, address. Click activates; AI cannot activate.", title_color=C_BLUE)
    add_card(s3, Inches(6.8), Inches(2.4), Inches(2.7), Inches(2.0),
             "3 · Verify & Reserve", "Code rehydrates catalog facts inside BEGIN IMMEDIATE transaction.", title_color=C_GREEN, border_color=C_GREEN)
    add_card(s3, Inches(9.8), Inches(2.4), Inches(2.7), Inches(2.0),
             "4 · Claim & Dispatch", "One owner claims Action Grant and dispatches Razorpay create_payment_link.", title_color=C_BLUE)

    add_card(s3, Inches(0.8), Inches(4.7), Inches(5.7), Inches(1.7),
             "The LLM never holds authority",
             "It translates intent and ranks substitutes. Server-owned SKUs and integer-paise prices determine the canonical cart facts.",
             title_color=C_TEXT)
    add_card(s3, Inches(6.8), Inches(4.7), Inches(5.7), Inches(1.7),
             "The actuator trusts verified state",
             "Envelope hash, policy revision, merchant ID, slot tags, and one-owner compare-and-set must all match before network dispatch.",
             title_color=C_TEXT)
    add_footer(s3, "POST /envelopes/draft → draft intent · POST /autopilot/execute → deterministic authorization gate")

    # --- SLIDE 4: The Differentiator ---
    s4 = prs.slides.add_slide(blank_layout)
    set_bg(s4)
    add_header(s4, "The differentiator", "An exact grant, not a broad permissions token.",
               "Each authorization binds who, what, which merchant, exact slot tags, and one purchase attempt. Any mismatch fails closed.")

    add_card(s4, Inches(0.8), Inches(2.4), Inches(5.7), Inches(3.8),
             "Bindings Enforced at Dispatch",
             "• Envelope identity: id · version · hash (immutable spend fence)\n"
             "• Cart + amount: cart_hash · quote_hash (integer paise only)\n"
             "• Context binding: user_id · merchant_id · fulfillment_profile\n"
             "• Single dispatch: purchase_attempt_id (CAS owner token)\n"
             "• Ambiguity-safe: UNKNOWN state holds headroom, blocks auto-retry\n"
             "• Cross-envelope ceiling: user aggregate exposure check",
             title_color=C_BLUE)

    add_card(s4, Inches(6.8), Inches(2.4), Inches(5.7), Inches(3.8),
             "In-Transaction Verification Gate",
             "authorize_and_reserve(request):\n\n"
             "  BEGIN IMMEDIATE\n"
             "  re-read envelope row, version & hash\n"
             "  check aggregate user authority ceiling\n"
             "  verify merchant, address, deadline, slots\n"
             "  if quote > cap: return PolicyDelta\n"
             "  mint exact Action Grant (one use)\n\n"
             "  # 8 concurrent attempts → exactly 1 winner",
             title_color=C_GREEN, border_color=C_GREEN)

    add_footer(s4, "Enforced atomically inside SQLite BEGIN IMMEDIATE before any network I/O")

    # --- SLIDE 5: Razorpay-Native Execution ---
    s5 = prs.slides.add_slide(blank_layout)
    set_bg(s5)
    add_header(s5, "Razorpay-native execution", "One registered action; dual HMAC-SHA256 evidence.",
               "Exposes create_payment_link through a Razorpay MCP-compatible transport. Receipts decouple authorization from lifecycle state.")

    add_card(s5, Inches(0.8), Inches(2.4), Inches(5.7), Inches(3.8),
             "Strict Actuator Boundary",
             "• canonicalize_action: Strict amount, currency, reference, notes, partial-payment flag, and schema hash.\n\n"
             "• claim_action_grant: Current policy, exact bindings, and one-owner compare-and-set before network I/O.\n\n"
             "• create_payment_link: The result is ACTION_ISSUED. Payment link is created; settlement requires provider evidence.",
             title_color=C_BLUE)

    add_card(s5, Inches(6.8), Inches(2.4), Inches(5.7), Inches(3.8),
             "Dual-Signature Action Receipt",
             "ReceiptAuthorization (immutable core):\n"
             "  → grant_id, cart_hash, quote_hash, policy_hash\n"
             "  → signed by authorization_signature\n\n"
             "ReceiptStatus (mutable lifecycle):\n"
             "  → state (action_issued / settled / unknown)\n"
             "  → signed by status_signature\n\n"
             "Auth core survives settlement verification cleanly.",
             title_color=C_GREEN, border_color=C_GREEN)

    add_footer(s5, "Vulcan is context, not a dependency: Razorpay positions it for checkout intelligence; Action Firewall gates authorization.")

    # --- SLIDE 6: Cross-Envelope Authority & Concurrency ---
    s6 = prs.slides.add_slide(blank_layout)
    set_bg(s6)
    add_header(s6, "Cross-envelope control", "A user authority ceiling spans multiple jobs.",
               "While individual Purchase Envelopes bound single purchases, an aggregate user ceiling prevents multi-agent spend accumulation.")

    add_card(s6, Inches(0.8), Inches(2.4), Inches(3.7), Inches(2.8),
             "Authority Ceiling",
             "Aggregate ceiling table (₹2,000 / weekly) bounds total outstanding exposure across all active envelopes.\n\nInspectable via GET /authority?user_id=.",
             title_color=C_BLUE)
    add_card(s6, Inches(4.8), Inches(2.4), Inches(3.7), Inches(2.8),
             "Atomic Concurrency",
             "N parallel threads colliding across separate envelopes are gated under BEGIN IMMEDIATE; exactly k succeed, remainder blocked with BLOCK_USER_CEILING_EXCEEDED.",
             title_color=C_GREEN, border_color=C_GREEN)
    add_card(s6, Inches(8.8), Inches(2.4), Inches(3.7), Inches(2.8),
             "Headroom Accounting",
             "Committed and pending exposure count against ceiling immediately. Unresolved exposure is held safely until authoritative settlement or reconciliation.",
             title_color=C_AMBER)

    add_card(s6, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.9),
             "Concurrency Proof",
             "8 concurrent threads under one envelope produce exactly 1 issued action; parallel envelopes collide at user ceiling with zero overspend.",
             title_color=C_TEXT, bg_color=C_PANEL2)
    add_footer(s6, "GET /authority?user_id= → real-time aggregated headroom and active envelope count")

    # --- SLIDE 7: Live Demo Beats ---
    s7 = prs.slides.add_slide(blank_layout)
    set_bg(s7)
    add_header(s7, "Live demo", "Three beats: useful autonomy, hard refusal, ambiguity.",
               "The demo is designed around one useful recovery, one field-specific refusal, and one distributed hold.")

    add_card(s7, Inches(0.8), Inches(2.4), Inches(3.7), Inches(2.8),
             "1 · Useful Autonomy",
             "Stock loss → safe substitute\n\n"
             "Preferred pasta sells out. System re-plans an eligible substitute and issues one simulated link without asking the shopper again.\n\n"
             "[ALLOW_ENVELOPE → ACTION_ISSUED]",
             title_color=C_GREEN, border_color=C_GREEN)
    add_card(s7, Inches(4.8), Inches(2.4), Inches(3.7), Inches(2.8),
             "2 · Hard Refusal",
             "Merchant changed → refuse\n\n"
             "Quote moves to unapproved merchant. Refused with field-level PolicyDelta (merchant_id). Zero actuator calls.\n\n"
             "[BLOCK_ENVELOPE_MISMATCH]",
             title_color=C_RED, border_color=C_RED)
    add_card(s7, Inches(8.8), Inches(2.4), Inches(3.7), Inches(2.8),
             "3 · Safe Ambiguity",
             "Provider timeout → hold\n\n"
             "Provider times out after dispatch. State held as UNKNOWN; envelope use remains occupied; duplicate retry suppressed.\n\n"
             "[UNKNOWN OUTCOME]",
             title_color=C_AMBER, border_color=C_AMBER)

    add_card(s7, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.9),
             "Fault Injection (Demo Only)",
             "Scenarios are injected by the operator to exercise deterministic boundaries on camera; endpoint refuses them outside DEMO_MODE.",
             title_color=C_TEXT, bg_color=C_PANEL2)
    add_footer(s7, "Terminal rehearsal: python scripts/demo_autopilot.py runs all beats offline with simulated provider")

    # --- SLIDE 8: Reproducible Proof ---
    s8 = prs.slides.add_slide(blank_layout)
    set_bg(s8)
    add_header(s8, "Reproducible proof", "Deterministic properties, verified on this machine.",
               "Synthetic authorization-correctness evidence, not inflated conversion claims.")

    add_card(s8, Inches(0.8), Inches(2.4), Inches(2.7), Inches(2.6),
             "146", "Backend tests passing\n\nExact suite: policy limits, envelopes, receipts, HTTP contracts, and concurrency.",
             title_color=C_GREEN, border_color=C_GREEN)
    add_card(s8, Inches(3.8), Inches(2.4), Inches(2.7), Inches(2.6),
             "650", "Generated boundary cases\n\n10 goal families, 104 distinct carts across 13 adversarial drift families (100% pass).",
             title_color=C_BLUE, border_color=C_BLUE)
    add_card(s8, Inches(6.8), Inches(2.4), Inches(2.7), Inches(2.6),
             "550 / 550", "Violations blocked\n\nPrice, merchant, category, destination, and tampering breaches fail closed.",
             title_color=C_RED, border_color=C_RED)
    add_card(s8, Inches(9.8), Inches(2.4), Inches(2.7), Inches(2.6),
             "8 → 1", "Concurrent claims to calls\n\nEight parallel attempts under one envelope; exactly one issued action.",
             title_color=C_GREEN, border_color=C_GREEN)

    add_card(s8, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.1),
             "Empirical Verification Scope Caveat",
             "Caveat: Synthetic deterministic authorization correctness. Not a claim of production conversion, GMV, fraud reduction, or payment success. Frontend clean; npm audit: 0; integer paise throughout.",
             title_color=C_TEXT, bg_color=C_PANEL2)
    add_footer(s8, "Reproduce in seconds: python -m pytest -q && python scripts/evaluate_autopilot.py")

    # --- SLIDE 9: Objections + Limits ---
    s9 = prs.slides.add_slide(blank_layout)
    set_bg(s9)
    add_header(s9, "Objections + limits", "What this proves — and what remains production work.",
               "Clear architectural boundaries and honest limitations.")

    add_card(s9, Inches(0.8), Inches(2.4), Inches(5.7), Inches(1.8),
             "“Rules engine with AI decoration?”",
             "AI translates open-ended goals and ranks substitutes. The payment boundary is deterministic because model creativity is not authorization.",
             title_color=C_BLUE)
    add_card(s9, Inches(6.8), Inches(2.4), Inches(5.7), Inches(1.8),
             "“Policy or envelope changes mid-checkout?”",
             "Version and hashes are revalidated at claim. A revoked envelope or stale grant is cancelled before network I/O.",
             title_color=C_BLUE)
    add_card(s9, Inches(0.8), Inches(4.5), Inches(5.7), Inches(1.8),
             "“Razorpay times out after accepting?”",
             "State becomes UNKNOWN, retains exposure, and rejects auto-retry. Reconciler applies authoritative provider events.",
             title_color=C_BLUE)
    add_card(s9, Inches(6.8), Inches(4.5), Inches(5.7), Inches(1.8),
             "“Production hardening path?”",
             "Single-instance SQLite serialization is an honest proof. Production requires authenticated principals, tenant scoping, durable outbox, and signed webhooks.",
             title_color=C_BLUE)
    add_footer(s9, "Honest limitations: application-signed evidence only; synthetic merchant environment; pull-based reconciliation")

    # --- SLIDE 10: Working Submission ---
    s10 = prs.slides.add_slide(blank_layout)
    set_bg(s10)
    add_header(s10, "Track 01 · working submission", "One approval for the job. Zero authority beyond it.",
               "Action Firewall turns Razorpay agentic checkout into a bounded, reviewable application control plane.")

    add_card(s10, Inches(0.8), Inches(2.8), Inches(11.7), Inches(2.2),
             "Verified System Properties",
             "• 146 passing backend tests covering pure spend policy, envelope drafts, concurrency, and reconciliation\n"
             "• 650 synthetic evaluation cases across 10 goal families (104 distinct carts, 100% boundary pass rate)\n"
             "• Dual HMAC-SHA256 Action Receipts preserving authorization proof across grant settlement\n"
             "• Cross-envelope User Authority Ceiling enforced atomically under SQLite BEGIN IMMEDIATE\n"
             "• Zero unauthorized actuator calls across all hostile scenarios and offline rehearsals",
             title_color=C_GREEN, border_color=C_GREEN)

    add_card(s10, Inches(0.8), Inches(5.3), Inches(11.7), Inches(0.9),
             "Submission Status",
             "pytest → 146 passed · 650 eval cases · npm audit → 0 · offline rehearsals → passed",
             title_color=C_BLUE, bg_color=C_PANEL2)

    add_footer(s10, "Track 01: AI Growth & Agentic Commerce", "Aryan Singh · Razorpay AI Buildathon 2026")

    output_path = Path(__file__).resolve().parent.parent.parent / "docs" / "Razorpay_Buildathon_Action_Firewall_Deck.pptx"
    prs.save(str(output_path))
    print(f"Saved: {output_path}")


if __name__ == "__main__":
    build_deck()
