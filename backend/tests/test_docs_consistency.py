"""Ensure documentation, pitch decks, and test metrics remain strictly consistent."""
from __future__ import annotations

import re
import subprocess
import sys
from pathlib import Path
import pytest


def test_docs_and_decks_match_pytest_collected_count():
    repo_root = Path(__file__).resolve().parent.parent.parent
    backend_dir = repo_root / "backend"
    readme_path = repo_root / "README.md"
    pitch_deck_path = repo_root / "docs" / "pitch-deck.html"
    demo_script_path = repo_root / "docs" / "SAFE_AUTOPILOT_DEMO.md"

    # 1. Collect total test count via pytest --collect-only
    res = subprocess.run(
        [sys.executable, "-m", "pytest", "--collect-only", "-q"],
        cwd=str(backend_dir),
        capture_output=True,
        text=True,
        check=True,
    )
    match_collect = re.search(r"(\d+)\s+tests collected", res.stdout)
    assert match_collect, f"Could not parse collected tests count from output:\n{res.stdout}"
    collected_count = int(match_collect.group(1))

    # 2. Verify README.md test count
    readme_text = readme_path.read_text(encoding="utf-8")
    readme_match = re.search(r"\*\*(\d+) passing backend tests\*\*", readme_text)
    assert readme_match, "README.md missing '**<N> passing backend tests**'"
    readme_count = int(readme_match.group(1))
    assert readme_count == collected_count, (
        f"README.md test count ({readme_count}) does not match pytest collected count ({collected_count})"
    )

    # 3. Verify docs/pitch-deck.html test count
    deck_text = pitch_deck_path.read_text(encoding="utf-8")
    stat_match = re.search(r'<div class="stat [^"]*">(\d+)</div>\s*<h3>Backend tests passing</h3>', deck_text)
    assert stat_match, "docs/pitch-deck.html missing '<div class=\"stat ...\">N</div><h3>Backend tests passing</h3>'"
    deck_stat_count = int(stat_match.group(1))
    assert deck_stat_count == collected_count, (
        f"pitch-deck.html stat count ({deck_stat_count}) does not match collected count ({collected_count})"
    )

    badge_match = re.search(r"pytest\s*→\s*(\d+)\s*passed", deck_text)
    assert badge_match, "docs/pitch-deck.html missing 'pytest → N passed'"
    deck_badge_count = int(badge_match.group(1))
    assert deck_badge_count == collected_count, (
        f"pitch-deck.html badge count ({deck_badge_count}) does not match collected count ({collected_count})"
    )

    footer_match = re.search(r"·\s*(\d+)\s*tests\s*·\s*offline demo", deck_text)
    assert footer_match, "docs/pitch-deck.html missing '· N tests · offline demo'"
    deck_footer_count = int(footer_match.group(1))
    assert deck_footer_count == collected_count, (
        f"pitch-deck.html footer count ({deck_footer_count}) does not match collected count ({collected_count})"
    )

    # 4. Verify docs/SAFE_AUTOPILOT_DEMO.md
    demo_text = demo_script_path.read_text(encoding="utf-8")
    demo_match = re.search(r"The (\d+)-test integration suite", demo_text)
    assert demo_match, "docs/SAFE_AUTOPILOT_DEMO.md missing 'The N-test integration suite'"
    demo_count = int(demo_match.group(1))
    assert demo_count == collected_count, (
        f"SAFE_AUTOPILOT_DEMO.md test count ({demo_count}) does not match collected count ({collected_count})"
    )

    # 5. Verify pitch deck content requirements
    assert "envelope" in deck_text.lower(), "pitch-deck.html must mention 'envelope'"
    assert "autopilot" in deck_text.lower(), "pitch-deck.html must mention 'autopilot'"
    assert "650" in deck_text, "pitch-deck.html must mention 650 generated cases"
    assert "104" in deck_text, "pitch-deck.html must mention 104 distinct carts"
    assert "61 tests" not in deck_text, "pitch-deck.html contains obsolete test count '61 tests'"
    assert "61 passed" not in deck_text, "pitch-deck.html contains obsolete test count '61 passed'"
    assert "51" not in stat_match.group(0), "pitch-deck.html contains obsolete test count '51'"

    # 6. Verify Do-Not-Say adherence in pitch-deck.html
    assert "powered by vulcan" not in deck_text.lower(), "Violates Do Not Say: 'powered by Vulcan'"
    assert "npci" not in deck_text.lower(), "Violates Do Not Say: NPCI mandate claim"
    assert "upi mandate" not in deck_text.lower(), "Violates Do Not Say: UPI mandate claim"
    assert "regulatory mandate" not in deck_text.lower(), "Violates Do Not Say: regulatory mandate claim"
    assert "ap2 compliance" not in deck_text.lower(), "Violates Do Not Say: AP2 compliance claim"
    assert "ap2 compliant" not in deck_text.lower(), "Violates Do Not Say: AP2 compliant claim"

    # 7. Verify PPTX deck consistency
    pptx_path = repo_root / "docs" / "Razorpay_Buildathon_Action_Firewall_Deck.pptx"
    assert pptx_path.exists(), "PPTX deck file must exist"
    try:
        from pptx import Presentation
        prs = Presentation(str(pptx_path))
        pptx_text = " ".join(
            shape.text
            for slide in prs.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )
        assert str(collected_count) in pptx_text, f"PPTX deck missing test count {collected_count}"
        assert "650" in pptx_text, "PPTX deck missing '650'"
        assert "104" in pptx_text, "PPTX deck missing '104'"
        assert "envelope" in pptx_text.lower(), "PPTX deck missing 'envelope'"
        assert "autopilot" in pptx_text.lower(), "PPTX deck missing 'autopilot'"
        assert "Aryan Singh" in pptx_text, "PPTX deck missing author 'Aryan Singh'"
        assert "61" not in pptx_text, "PPTX deck contains obsolete test count '61'"
        assert "51" not in pptx_text, "PPTX deck contains obsolete test count '51'"
        assert "powered by vulcan" not in pptx_text.lower(), "PPTX deck violates Do Not Say: 'powered by Vulcan'"
    except ImportError:
        pass  # python-pptx optional if run in minimal test env

